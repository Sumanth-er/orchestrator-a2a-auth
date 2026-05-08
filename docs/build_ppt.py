"""Generate the demo-a2a architecture deck.

Run:  python docs/build_ppt.py
Out:  docs/demo-a2a-architecture.pptx
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- palette (simple, calm) -----------------------------------------
NAVY      = RGBColor(0x1F, 0x3A, 0x5F)
TEAL      = RGBColor(0x2E, 0x86, 0xAB)
SKY       = RGBColor(0xA7, 0xC9, 0xE3)
SAND      = RGBColor(0xF5, 0xEC, 0xD8)
GRAY      = RGBColor(0x6C, 0x75, 0x7D)
LIGHT     = RGBColor(0xF5, 0xF7, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x2E, 0x8B, 0x57)
RED       = RGBColor(0xC0, 0x39, 0x2B)
ORANGE    = RGBColor(0xE6, 0x8A, 0x00)
DARKTEXT  = RGBColor(0x1F, 0x29, 0x37)


# ---------- helpers ---------------------------------------------------------
def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    if not shadow:
        shp.shadow.inherit = False
    shp.adjustments[0] = 0.08
    return shp


def add_text(slide, x, y, w, h, text,
             *, size=14, bold=False, color=DARKTEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARKTEXT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, weight=2.0, label=None,
              label_above=True):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # add an end-arrow via XML
    ln = line.line._get_or_add_ln()  # type: ignore[attr-defined]
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    if label:
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        offy = Emu(-220000) if label_above else Emu(40000)
        tb = slide.shapes.add_textbox(mid_x - Inches(1.2), mid_y + offy,
                                      Inches(2.4), Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(10)
        r.font.italic = True
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return line


def slide_title(slide, title, subtitle=None):
    # top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 Inches(13.333), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, Inches(0.4), Inches(0.05), Inches(12.5), Inches(0.5),
             title, size=22, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.4),
                 subtitle, size=13, color=GRAY, anchor=MSO_ANCHOR.TOP)


def footer(slide, page_no):
    add_text(slide, Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.3),
             f"demo-a2a  ·  A2A + Keycloak auth  ·  slide {page_no}",
             size=9, color=GRAY)


# ---------- slide builders --------------------------------------------------
def slide_cover(prs, n):
    s = add_blank(prs)
    # full-bleed colored panel
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0),
                               Inches(13.333), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()

    add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1.2),
             "demo-a2a", size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.0), Inches(12), Inches(0.8),
             "A2A Protocol + Keycloak Authentication & Authorization",
             size=24, color=SKY)
    add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.5),
             "Architecture & End-to-End Flow",
             size=18, color=SAND)

    # tag chips
    tags = ["OAuth2 / OIDC", "PKCE", "RFC 8693 Token-Exchange", "JWT", "SSE", "FastAPI"]
    x = Inches(0.8); y = Inches(5.2)
    for t in tags:
        chip = add_rect(s, x, y, Inches(2.1), Inches(0.5), TEAL)
        add_text(s, x, y, Inches(2.1), Inches(0.5), t, size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(2.2)


def slide_overview(prs, n):
    s = add_blank(prs)
    slide_title(s, "Project Overview",
                "What demo-a2a demonstrates and why it exists")

    # left column: goals
    add_text(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.4),
             "Goals", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(4.5), [
        "Browser → orchestrator → multiple A2A agents, all behind Keycloak.",
        "Real OAuth2 / OIDC PKCE login (no fake bearer).",
        "RFC 8693 Token-Exchange between hops — each agent gets a token whose 'aud' matches it.",
        "Per-agent authorization at the agent itself (audience + realm role).",
        "Graceful denial: never raises 403 to the user — UI shows a red 'no access' chip.",
        "Drop-in agent registration: copy a folder, add one registry line.",
    ], size=14)

    # right column: components
    add_text(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4),
             "Components", size=18, bold=True, color=NAVY)
    comps = [
        ("Browser UI",       "static HTML + oidc-client-ts (PKCE)"),
        ("Orchestrator",     "FastAPI · create_agent runtime · SkillMiddleware · call_agent tool"),
        ("Weather Agent",    "A2A server · LangChain + Ollama"),
        ("Billing Agent",    "A2A server · LangChain + Ollama"),
        ("Keycloak",         "Realm 'a2a-demo' · clients · roles · audience scopes"),
        ("Shared lib",       "JwtValidator · TokenExchanger · KeycloakAuthMiddleware"),
    ]
    y = Inches(1.6)
    for name, desc in comps:
        add_rect(s, Inches(7.0), y, Inches(2.3), Inches(0.55), TEAL)
        add_text(s, Inches(7.0), y, Inches(2.3), Inches(0.55), name,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.4), y, Inches(3.8), Inches(0.55), desc,
                 size=12, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.7)

    footer(s, n)


def slide_architecture(prs, n):
    s = add_blank(prs)
    slide_title(s, "High-Level Architecture",
                "Five processes, one realm, two policy checkpoints")

    # Browser
    add_rect(s, Inches(0.4), Inches(2.6), Inches(2.0), Inches(1.2), SAND)
    add_text(s, Inches(0.4), Inches(2.6), Inches(2.0), Inches(0.4),
             "Browser", size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.4), Inches(2.95), Inches(2.0), Inches(0.85),
             "static UI\noidc-client-ts\nPKCE login",
             size=11, color=DARKTEXT, align=PP_ALIGN.CENTER)

    # Orchestrator
    add_rect(s, Inches(3.2), Inches(2.2), Inches(3.0), Inches(2.0), SKY)
    add_text(s, Inches(3.2), Inches(2.2), Inches(3.0), Inches(0.45),
             "Orchestrator :3000", size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.2), Inches(2.6), Inches(3.0), Inches(1.6),
             "FastAPI / SSE\nfront-door JwtValidator\nLLM router (Ollama)\nTokenExchanger\nA2A Dispatcher",
             size=11, color=DARKTEXT, align=PP_ALIGN.CENTER)

    # Agents stack
    add_rect(s, Inches(7.0), Inches(1.5), Inches(3.0), Inches(1.4), LIGHT, line=TEAL)
    add_text(s, Inches(7.0), Inches(1.5), Inches(3.0), Inches(0.45),
             "Weather Agent :9101", size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.0), Inches(1.95), Inches(3.0), Inches(0.95),
             "A2A server\nKeycloakAuthMiddleware\naud=weather-agent\nrole=weather.read",
             size=10, color=DARKTEXT, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(7.0), Inches(3.1), Inches(3.0), Inches(1.4), LIGHT, line=TEAL)
    add_text(s, Inches(7.0), Inches(3.1), Inches(3.0), Inches(0.45),
             "Billing Agent :9102", size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.0), Inches(3.55), Inches(3.0), Inches(0.95),
             "A2A server\nKeycloakAuthMiddleware\naud=billing-agent\nrole=billing.read",
             size=10, color=DARKTEXT, align=PP_ALIGN.CENTER)

    # Keycloak
    add_rect(s, Inches(10.6), Inches(2.4), Inches(2.5), Inches(1.6), NAVY)
    add_text(s, Inches(10.6), Inches(2.4), Inches(2.5), Inches(0.45),
             "Keycloak :8080", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.6), Inches(2.85), Inches(2.5), Inches(1.15),
             "realm a2a-demo\nclients · roles\naud-* scopes\nJWKS · token-exchange",
             size=11, color=SAND, align=PP_ALIGN.CENTER)

    # Arrows
    add_arrow(s, Inches(2.4), Inches(3.2), Inches(3.2), Inches(3.2),
              color=NAVY, label="Bearer (aud=orchestrator)")
    add_arrow(s, Inches(6.2), Inches(2.6), Inches(7.0), Inches(2.2),
              color=GREEN, label="Bearer (aud=weather)")
    add_arrow(s, Inches(6.2), Inches(3.6), Inches(7.0), Inches(3.7),
              color=GREEN, label="Bearer (aud=billing)")
    add_arrow(s, Inches(6.2), Inches(2.95), Inches(10.6), Inches(2.85),
              color=ORANGE, label="token-exchange (RFC 8693)", label_above=False)
    add_arrow(s, Inches(4.7), Inches(4.2), Inches(11.0), Inches(4.0),
              color=GRAY, weight=1.0, label="JWKS verify",
              label_above=False)

    # legend
    add_text(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.4),
             "Legend",
             size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.4), [
        "Blue arrow — user JWT (aud=orchestrator). Front-door auth only.",
        "Orange arrow — RFC 8693 Token-Exchange against Keycloak (orchestrator → aud=<agent>).",
        "Green arrow — exchanged JWT carried to the agent; agent validates aud + role.",
        "Gray arrow — JWKS public-key fetch; cached 10 min.",
    ], size=12)

    footer(s, n)


def slide_keycloak(prs, n):
    s = add_blank(prs)
    slide_title(s, "Keycloak Realm Layout",
                "Clients, scopes, and the audience-claim chain")

    add_text(s, Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.4),
             "Realm: a2a-demo", size=16, bold=True, color=NAVY)

    # 4 client cards
    clients = [
        ("a2a-ui",        "Public · PKCE",          "browser login\nscope: openid profile",       SAND),
        ("orchestrator",  "Confidential",           "token-exchange initiator\nscopes: aud-weather, aud-billing", SKY),
        ("weather-agent", "Bearer-only",            "validates JWTs\nrequires aud=weather-agent",                 LIGHT),
        ("billing-agent", "Bearer-only",            "validates JWTs\nrequires aud=billing-agent",                 LIGHT),
    ]
    x = Inches(0.5)
    for name, kind, desc, color in clients:
        add_rect(s, x, Inches(1.7), Inches(3.0), Inches(2.0), color, line=TEAL)
        add_text(s, x, Inches(1.75), Inches(3.0), Inches(0.4), name,
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.15), Inches(3.0), Inches(0.35), kind,
                 size=11, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.55), Inches(3.0), Inches(1.1), desc,
                 size=11, color=DARKTEXT, align=PP_ALIGN.CENTER)
        x += Inches(3.2)

    # roles & users
    add_text(s, Inches(0.5), Inches(4.0), Inches(6.0), Inches(0.4),
             "Realm roles", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(4.4), Inches(6.0), Inches(1.5), [
        "weather.read — required by weather-agent middleware",
        "billing.read — required by billing-agent middleware",
        "Roles are checked at the agent, not the orchestrator",
    ], size=12)

    add_text(s, Inches(7.0), Inches(4.0), Inches(6.0), Inches(0.4),
             "Test users", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(4.4), Inches(6.0), Inches(1.5), [
        "alice — weather.read + billing.read  → both agents work",
        "bob   — weather.read only            → billing returns 403",
        "Add roles via Users → Role mapping",
    ], size=12)

    # audience scope explanation
    add_rect(s, Inches(0.5), Inches(5.9), Inches(12.5), Inches(1.1), SAND)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.4),
             "How the audience claim flows",
             size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.7),
             "1) UI logs in → token aud=orchestrator    "
             "2) Orchestrator does token-exchange → token aud=weather-agent (or billing-agent)    "
             "3) Agent JwtValidator enforces aud=<itself>",
             size=12, color=DARKTEXT)

    footer(s, n)


def slide_hops(prs, n):
    s = add_blank(prs)
    slide_title(s, "Eight Hops of a Single Chat Message",
                "Where each auth check fires, in order — including the new agent runtime")

    headers = ["#", "Hop", "File · Function", "Check enforced"]
    rows = [
        ("0", "Browser PKCE login",       "static/app.js · init / signIn",
         "Keycloak verifies code_verifier + creds"),
        ("1", "Browser → orchestrator",    "static/app.js · sendMessage:145–161",
         "Authorization: Bearer attached"),
        ("2", "Front-door auth at /chat",  "orchestrator/main.py · _authenticate:73–80\n+ shared/a2a_auth/jwt_validator.py · validate",
         "signature + iss + exp (no aud)"),
        ("3", "Build per-request runtime", "orchestrator/agent_runtime.py · build_runtime\n+ orchestrator/agent_tools.py · make_tool",
         "user_token closed over · auth-blind middleware"),
        ("4", "LLM emits N tool_calls",     "orchestrator/skills_middleware.py · wrap_model_call\n+ create_agent ReAct loop",
         "system prompt = base + skills addendum"),
        ("5", "RFC 8693 token-exchange",   "orchestrator/a2a_dispatcher.py · call:60\n+ shared/a2a_auth/token_exchange.py · exchange",
         "Keycloak grants aud=<agent> token"),
        ("6", "Orchestrator → agent",      "orchestrator/a2a_dispatcher.py · _invoke:99–122\n+ shared/a2a_auth/client_auth.py · bearer_httpx_client",
         "exchanged JWT carried downstream"),
        ("7", "Agent middleware",          "shared/a2a_auth/server_middleware.py · dispatch",
         "sig + iss + exp + aud + realm role"),
        ("8", "Failure surfacing → UI",    "a2a_dispatcher.py · _map_http_error\n→ static/app.js · addChip",
         "401/403 → red 'no access' chip"),
    ]

    # table area
    top = Inches(1.2); left = Inches(0.4)
    col_w = [Inches(0.4), Inches(2.5), Inches(5.6), Inches(4.4)]
    row_h = Inches(0.62)

    # header row
    x = left
    for i, h in enumerate(headers):
        add_rect(s, x, top, col_w[i], Inches(0.4), NAVY)
        add_text(s, x, top, col_w[i], Inches(0.4), h, size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    # body rows
    y = top + Inches(0.4)
    for i, row in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        x = left
        for j, cell in enumerate(row):
            add_rect(s, x, y, col_w[j], row_h, bg, line=GRAY)
            sz = 10 if j in (2, 3) else 11
            bold = (j == 0)
            color = NAVY if j == 0 else DARKTEXT
            add_text(s, x, y, col_w[j], row_h, cell, size=sz, bold=bold,
                     color=color, align=PP_ALIGN.LEFT if j > 0 else PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[j]
        y += row_h

    footer(s, n)


def slide_flow_diagram(prs, n):
    s = add_blank(prs)
    slide_title(s, "End-to-End Flow Diagram",
                "Sequence across browser, orchestrator, Keycloak, agent")

    # vertical lanes
    lanes = [
        ("Browser",      Inches(0.7),  SAND),
        ("Orchestrator", Inches(3.6),  SKY),
        ("Keycloak",     Inches(6.7),  NAVY),
        ("Agent",        Inches(10.0), LIGHT),
    ]
    for label, x, color in lanes:
        add_rect(s, x, Inches(1.2), Inches(2.6), Inches(0.5), color)
        text_color = WHITE if color in (NAVY,) else NAVY
        add_text(s, x, Inches(1.2), Inches(2.6), Inches(0.5), label,
                 size=13, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # lane line
        ln = s.shapes.add_connector(1, x + Inches(1.3), Inches(1.7),
                                    x + Inches(1.3), Inches(6.7))
        ln.line.color.rgb = GRAY
        ln.line.width = Pt(0.75)
        ln.line.dash_style = 7  # dash

    # message arrows
    msgs = [
        # (from_x, to_x, y, label, color)
        (2.0, 4.9,  1.95, "1. PKCE redirect → login", NAVY),
        (4.9, 2.0,  2.30, "2. ?code=…  (callback)",   NAVY),
        (2.0, 4.9,  2.70, "3. POST /chat + Bearer (aud=orchestrator)", TEAL),
        (4.9, 4.9,  3.05, "4. _authenticate → build_runtime (LLM picks N agents)", GRAY),
        (4.9, 8.0,  3.45, "5. token-exchange × N (aud=weather, aud=billing, …)", ORANGE),
        (8.0, 4.9,  3.85, "6. exchanged JWTs", ORANGE),
        (4.9, 11.3, 4.25, "7. parallel JSON-RPC × N (one per tool_call)", GREEN),
        (11.3, 11.3, 4.65, "8. each agent: aud + role check", GRAY),
        (11.3, 4.9, 5.05, "9. results merged in summarize-skill", GREEN),
        (4.9, 2.0,  5.45, "10. SSE: N× (agent_selected, agent_result) → reply", TEAL),
    ]
    for fx, tx, y, label, color in msgs:
        add_arrow(s, Inches(fx), Inches(y), Inches(tx), Inches(y),
                  color=color, weight=1.6, label=label, label_above=True)

    # bottom note
    add_rect(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.55), SAND)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.4), Inches(0.55),
             "Two policy checkpoints: step 5 (Keycloak decides aud) and step 8 "
             "(Agent decides role). Steps 5–8 fan out per tool_call — one user prompt → 1..N agents in parallel.",
             size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n)


def slide_orchestrator(prs, n):
    s = add_blank(prs)
    slide_title(s, "Orchestrator — Internals",
                "FastAPI app, lifespan-built singletons, agent runtime, SSE streaming")

    # left: file tree
    add_text(s, Inches(0.4), Inches(1.2), Inches(5), Inches(0.4),
             "orchestrator/ files", size=14, bold=True, color=NAVY)
    files = [
        ("main.py",              "/chat SSE · _authenticate · event-bus"),
        ("agent_runtime.py",     "build_runtime · run · ChatOllama"),
        ("agent_tools.py",       "make_tool — single call_agent tool"),
        ("skills_middleware.py", "SkillMiddleware · load_skill · SKILLS"),
        ("a2a_dispatcher.py",    "Dispatcher · _invoke · _map_http_error"),
        ("registry.py",          "AgentEntry list (name, url, audience, desc)"),
        ("static/app.js",        "PKCE · sendMessage · parseSse · chips"),
    ]
    y = Inches(1.6)
    for name, desc in files:
        add_rect(s, Inches(0.4), y, Inches(2.5), Inches(0.45), TEAL)
        add_text(s, Inches(0.4), y, Inches(2.5), Inches(0.45), name,
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas")
        add_text(s, Inches(3.0), y, Inches(3.6), Inches(0.45), desc,
                 size=10, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.5)

    # right: lifespan
    add_text(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             "Singletons (built in main.py · lifespan)",
             size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(1.6), Inches(6), Inches(2.6), [
        "KeycloakSettings — issuer, jwks_uri, token_endpoint",
        "JwtValidator — JWKS-cached PyJWT decoder",
        "TokenExchanger — RFC 8693 client w/ (jti, aud) cache",
        "Dispatcher — wraps exchanger, calls A2A clients",
    ], size=12)

    add_text(s, Inches(7.0), Inches(4.3), Inches(6), Inches(0.4),
             "/chat SSE event stream (UNCHANGED — UI stays the same)",
             size=13, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(4.7), Inches(6), Inches(2.2), [
        "user_authenticated  → username chip",
        "agent_selected       → pending chip (now 1..N per turn)",
        "agent_result         → green ok / red denied / orange error",
        "no_agent             → small-talk path (no tool was called)",
        "reply                → final assistant bubble (summarize-skill)",
    ], size=12)

    footer(s, n)


def slide_runtime(prs, n):
    """NEW slide — explains the create_agent runtime + SkillMiddleware."""
    s = add_blank(prs)
    slide_title(s, "Agent Runtime + SkillMiddleware",
                "How one user prompt becomes 1..N parallel agent calls")

    # left: 3 box stack — base prompt / middleware / tool
    boxes = [
        ("BASE_SYSTEM",
         "agent_runtime.py · BASE_SYSTEM\n\n"
         "“You are a coordinator… delegate to specialist agents via "
         "call_agent… call them in parallel when independent.”",
         SKY),
        ("SkillMiddleware.wrap_model_call",
         "skills_middleware.py · SkillMiddleware\n\n"
         "Appends a '## Available Skills' addendum to system_message "
         "via request.override(...). Owns load_skill tool as class var "
         "(tools = [load_skill]).",
         SAND),
        ("call_agent tool",
         "agent_tools.py · make_tool\n\n"
         "Args: agent (Literal name), query (rephrased sub-question). "
         "Closes over user_token + dispatcher + event_sink. "
         "Emits agent_selected / agent_result events.",
         LIGHT),
    ]
    y = Inches(1.3)
    for title, body, color in boxes:
        add_rect(s, Inches(0.4), y, Inches(6.0), Inches(1.6), color, line=TEAL)
        add_text(s, Inches(0.4), y + Inches(0.05), Inches(6.0), Inches(0.4),
                 title, size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.55), y + Inches(0.45), Inches(5.7), Inches(1.1),
                 body, size=10, color=DARKTEXT)
        y += Inches(1.75)

    # right: data flow / SKILLS catalog
    add_text(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(0.4),
             "Skill catalog (skills_middleware.py · SKILLS)",
             size=13, bold=True, color=NAVY)
    skills = [
        ("plan",       "decompose hard prompts into sub-tasks"),
        ("execute",    "call call_agent (parallel for independent tasks)"),
        ("summarize",  "fuse all tool results into one reply"),
        ("smalltalk",  "answer directly without any tool call"),
    ]
    y = Inches(1.7)
    for name, desc in skills:
        add_rect(s, Inches(6.8), y, Inches(1.5), Inches(0.45), TEAL)
        add_text(s, Inches(6.8), y, Inches(1.5), Inches(0.45), name,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas")
        add_text(s, Inches(8.4), y, Inches(4.6), Inches(0.45), desc,
                 size=11, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.55)

    add_text(s, Inches(6.8), Inches(4.2), Inches(6.2), Inches(0.4),
             "Auth posture (UNCHANGED)",
             size=13, bold=True, color=NAVY)
    add_bullets(s, Inches(6.8), Inches(4.6), Inches(6.2), Inches(2.4), [
        "user_token NEVER enters graph state — captured by tool closure only",
        "Middleware reads/writes system_message — auth-blind by construction",
        "Dispatcher.call(entry, user_token, query) — same code path as before",
        "Per-tool-call: token-exchange → aud=<agent> → role check at agent",
        "SSE event names unchanged → static/app.js needs zero changes",
    ], size=11)

    footer(s, n)


def slide_agent_internals(prs, n):
    s = add_blank(prs)
    slide_title(s, "Agent — Internals",
                "Reusable factory + drop-in middleware + custom executor")

    # boxes top row
    boxes = [
        ("base_agent.py · make_agent_app",
         "Builds Starlette app with:\n• create_agent_card_routes(card)\n• create_jsonrpc_routes(handler, '/')\n• DefaultRequestHandler + InMemoryTaskStore\n• KeycloakAuthMiddleware",
         SKY),
        ("server_middleware.py · KeycloakAuthMiddleware",
         "On every non-public request:\n• extract_bearer\n• validate(token, expected_audience)\n• extract_realm_roles → require role\n• 401/403 JSONResponse on failure",
         SAND),
        ("agents/<name>/executor.py",
         "Per-agent business logic:\n• reads context.get_user_input()\n• creates new_task / artifacts\n• optional fine-grained role gating\n• knows nothing about JWTs",
         LIGHT),
    ]
    x = Inches(0.4)
    for title, body, color in boxes:
        add_rect(s, x, Inches(1.3), Inches(4.1), Inches(2.6), color, line=TEAL)
        add_text(s, x, Inches(1.4), Inches(4.1), Inches(0.5), title,
                 size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, Inches(1.95), Inches(4.1), Inches(1.9), body,
                 size=11, color=DARKTEXT)
        x += Inches(4.3)

    # public paths note
    add_rect(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.9), SAND)
    add_text(s, Inches(0.6), Inches(4.25), Inches(12.2), Inches(0.4),
             "Public paths (no auth)",
             size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(4.6), Inches(12.2), Inches(0.5),
             "/.well-known/agent-card.json   ·   /.well-known/agent.json   ·   /agent-card",
             size=12, color=DARKTEXT, font="Consolas")

    # status outcomes
    add_text(s, Inches(0.4), Inches(5.3), Inches(12), Inches(0.4),
             "What the agent returns",
             size=14, bold=True, color=NAVY)
    chips = [
        ("200 OK",            "executor produced a TextArtifact",          GREEN),
        ("401 UNAUTHORIZED",  "missing/invalid token, wrong audience",     RED),
        ("403 ACCESS_DENIED", "valid token, role not granted",             ORANGE),
    ]
    x = Inches(0.4); y = Inches(5.7)
    for code, why, color in chips:
        add_rect(s, x, y, Inches(4.1), Inches(1.0), color)
        add_text(s, x, y, Inches(4.1), Inches(0.45), code,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.45), Inches(4.1), Inches(0.55), why,
                 size=11, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += Inches(4.3)

    footer(s, n)


def slide_token_exchange(prs, n):
    s = add_blank(prs)
    slide_title(s, "RFC 8693 Token-Exchange · Detail",
                "How orchestrator obtains an aud=<agent> JWT")

    # request box
    add_rect(s, Inches(0.4), Inches(1.3), Inches(6.0), Inches(3.0), LIGHT, line=NAVY)
    add_text(s, Inches(0.4), Inches(1.35), Inches(6.0), Inches(0.4),
             "POST /realms/a2a-demo/protocol/openid-connect/token",
             size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    body = (
        "grant_type        = urn:ietf:params:oauth:grant-type:token-exchange\n"
        "client_id         = orchestrator\n"
        "client_secret     = ********\n"
        "subject_token     = <user JWT, aud=orchestrator>\n"
        "subject_token_type= urn:ietf:params:oauth:token-type:access_token\n"
        "requested_token_type = urn:ietf:params:oauth:token-type:access_token\n"
        "audience          = weather-agent"
    )
    add_text(s, Inches(0.55), Inches(1.8), Inches(5.8), Inches(2.5), body,
             size=11, color=DARKTEXT, font="Consolas")

    # response box
    add_rect(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(3.0), SAND, line=NAVY)
    add_text(s, Inches(7.0), Inches(1.35), Inches(6.0), Inches(0.4),
             "200 OK · response",
             size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    resp = (
        '{\n'
        '  "access_token": "<JWT>",   ← aud=weather-agent\n'
        '  "token_type":   "Bearer",\n'
        '  "expires_in":   300,\n'
        '  "issued_token_type":\n'
        '    "urn:ietf:params:oauth:token-type:access_token"\n'
        '}'
    )
    add_text(s, Inches(7.15), Inches(1.8), Inches(5.7), Inches(2.5), resp,
             size=11, color=DARKTEXT, font="Consolas")

    # caching
    add_text(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.4),
             "Caching (token_exchange.py · TokenExchanger)",
             size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(2), [
        "Cache key: (jti, audience). When user token refreshes, jti changes → fresh exchange.",
        "Refresh skew: 30 s — tokens are reused until <30 s remaining.",
        "Per-key asyncio.Lock prevents thundering herd on first miss.",
        "Failures from Keycloak → TokenExchangeError(403) → dispatcher returns AgentCallResult(status='denied').",
    ], size=13)

    footer(s, n)


def slide_failure_paths(prs, n):
    s = add_blank(prs)
    slide_title(s, "Failure Paths & UI Chips",
                "Every error becomes a chip — never a 500 to the user")

    # flowchart-ish
    boxes = [
        ("dispatcher.call(entry, user_token)", Inches(4.5), Inches(1.3), SKY,
         "orchestrator/a2a_dispatcher.py:51–79"),
        ("token-exchange OK?",               Inches(4.5), Inches(2.3), SAND, ""),
        ("call agent _invoke",               Inches(4.5), Inches(3.3), SKY,
         "a2a_dispatcher.py:99–122"),
        ("HTTP status?",                      Inches(4.5), Inches(4.3), SAND, ""),
    ]
    for title, x, y, color, sub in boxes:
        add_rect(s, x, y, Inches(4.0), Inches(0.7), color, line=NAVY)
        add_text(s, x, y, Inches(4.0), Inches(0.7), title,
                 size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            add_text(s, x, y + Inches(0.7), Inches(4.0), Inches(0.3), sub,
                     size=9, color=GRAY, align=PP_ALIGN.CENTER, font="Consolas")

    # arrows down
    add_arrow(s, Inches(6.5), Inches(2.0), Inches(6.5), Inches(2.3), color=NAVY)
    add_arrow(s, Inches(6.5), Inches(3.0), Inches(6.5), Inches(3.3), color=NAVY,
              label="yes")
    add_arrow(s, Inches(6.5), Inches(4.0), Inches(6.5), Inches(4.3), color=NAVY)

    # outcome chips bottom
    outcomes = [
        ("DENIED chip\nTokenExchangeError\n(no aud-* on orchestrator client)", RED, Inches(0.4), Inches(5.5)),
        ("DENIED chip\nagent 401 / 403\n(wrong aud or role)",                   RED, Inches(4.7), Inches(5.5)),
        ("ERROR chip\n5xx / unreachable\n(service down)",                       ORANGE, Inches(9.0), Inches(5.5)),
        ("OK chip\n200 + text artifact\n(executor ran)",                        GREEN,  Inches(0.4), Inches(6.5)),
    ]
    for body, color, x, y in outcomes:
        add_rect(s, x, y, Inches(3.9), Inches(0.9), color)
        add_text(s, x, y, Inches(3.9), Inches(0.9), body,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrow from token-exchange to denied
    add_arrow(s, Inches(4.5), Inches(2.65), Inches(2.4), Inches(5.5),
              color=RED, weight=1.4, label="no")

    # arrow from HTTP status to outcomes
    add_arrow(s, Inches(5.5), Inches(5.0), Inches(2.4), Inches(5.5),
              color=RED, weight=1.4, label="401")
    add_arrow(s, Inches(6.5), Inches(5.0), Inches(6.6), Inches(5.5),
              color=RED, weight=1.4, label="403")
    add_arrow(s, Inches(7.5), Inches(5.0), Inches(11.0), Inches(5.5),
              color=ORANGE, weight=1.4, label="5xx / other")
    add_arrow(s, Inches(5.5), Inches(5.0), Inches(2.4), Inches(6.5),
              color=GREEN, weight=1.4, label="200")

    footer(s, n)


def slide_add_agent(prs, n):
    s = add_blank(prs)
    slide_title(s, "Adding a New Agent",
                "Six steps · works for any LangChain / A2A executor")

    rows = [
        ("1", "Keycloak",
         "Create realm role <name>.read · client <name>-agent (bearer-only) · "
         "client scope aud-<name> with Audience mapper · assign aud-<name> "
         "as Default to orchestrator"),
        ("2", "Folder",
         "agents/<name>_agent/ with card.py · executor.py · main.py "
         "(copy from billing or weather)"),
        ("3", "Card",
         "AgentCard with id, name, description, examples — used by the "
         "router LLM to choose this agent"),
        ("4", "main.py",
         "make_agent_app(card=CARD, executor=YourExecutor(), "
         "expected_audience=<name>-agent, required_roles=[\"<name>.read\"])"),
        ("5", "Registry",
         "Append AgentEntry(name, url, audience, description) to "
         "orchestrator/registry.py"),
        ("6", "Launcher",
         "Add Service entry to run.py · assign role to test users"),
    ]
    top = Inches(1.2)
    for i, (num, title, body) in enumerate(rows):
        y = top + Inches(0.85 * i)
        add_rect(s, Inches(0.4), y, Inches(0.7), Inches(0.75), NAVY)
        add_text(s, Inches(0.4), y, Inches(0.7), Inches(0.75), num,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.2), y, Inches(2.0), Inches(0.75), TEAL)
        add_text(s, Inches(1.2), y, Inches(2.0), Inches(0.75), title,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.3), y, Inches(9.6), Inches(0.75), LIGHT, line=GRAY)
        add_text(s, Inches(3.3), y, Inches(9.6), Inches(0.75), body,
                 size=11, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n)


def slide_techstack(prs, n):
    s = add_blank(prs)
    slide_title(s, "Tech Stack & Files",
                "Quick reference for what's where")

    # left: stack
    add_text(s, Inches(0.4), Inches(1.2), Inches(6), Inches(0.4),
             "Stack", size=16, bold=True, color=NAVY)
    rows = [
        ("Python",   ">= 3.11"),
        ("Server",   "FastAPI · Starlette · uvicorn"),
        ("A2A",      "a2a-sdk[http-server] >= 0.3.0"),
        ("Runtime",  "langchain >=1.0 · langgraph · create_agent"),
        ("LLM",      "langchain-ollama (local Ollama models)"),
        ("Auth",     "PyJWT[crypto] · jwks · pydantic-settings"),
        ("UI",       "static HTML · oidc-client-ts · SSE"),
        ("Realm",    "Keycloak 25.x · token-exchange-standard feature"),
    ]
    y = Inches(1.6)
    for k, v in rows:
        add_rect(s, Inches(0.4), y, Inches(1.6), Inches(0.45), TEAL)
        add_text(s, Inches(0.4), y, Inches(1.6), Inches(0.45), k,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.1), y, Inches(4.5), Inches(0.45), v,
                 size=11, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.55)

    # right: file map
    add_text(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             "Key files", size=16, bold=True, color=NAVY)
    files = [
        ("orchestrator/main.py",                  "/chat SSE · _authenticate · event-bus"),
        ("orchestrator/agent_runtime.py",         "build_runtime · run · ChatOllama"),
        ("orchestrator/agent_tools.py",           "make_tool — single call_agent tool"),
        ("orchestrator/skills_middleware.py",     "SkillMiddleware · load_skill · SKILLS"),
        ("orchestrator/a2a_dispatcher.py",        "Dispatcher · _map_http_error"),
        ("orchestrator/registry.py",              "AgentEntry list"),
        ("orchestrator/static/app.js",            "PKCE · SSE · chips"),
        ("agents/base_agent.py",                  "make_agent_app factory"),
        ("agents/<name>/main.py",                 "wires audience + roles"),
        ("shared/a2a_auth/jwt_validator.py",      "JWKS + PyJWT decoder"),
        ("shared/a2a_auth/token_exchange.py",     "RFC 8693 client + cache"),
        ("shared/a2a_auth/server_middleware.py",  "agent-side auth"),
        ("shared/a2a_auth/client_auth.py",        "bearer_httpx_client hook"),
    ]
    y = Inches(1.6)
    for path, desc in files:
        add_text(s, Inches(7.0), y, Inches(3.6), Inches(0.32), path,
                 size=10, color=NAVY, font="Consolas")
        add_text(s, Inches(10.6), y, Inches(2.6), Inches(0.32), desc,
                 size=10, color=GRAY)
        y += Inches(0.36)

    footer(s, n)


def slide_summary(prs, n):
    s = add_blank(prs)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0),
                               Inches(13.333), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()

    add_text(s, Inches(0.8), Inches(0.8), Inches(12), Inches(0.8),
             "Summary", size=40, bold=True, color=WHITE)

    points = [
        "User identity is established once at the front door (Keycloak PKCE).",
        "Authorization is split: Keycloak controls who gets a token for which agent (audience).",
        "The agent itself controls who can run it (realm role + audience claim check).",
        "One user prompt → 1..N parallel agent calls via create_agent + a single call_agent tool.",
        "SkillMiddleware swaps the system prompt per stage — auth-blind by construction.",
        "All denials are graceful: red chips + LLM-phrased reply, never 500s.",
        "RFC 8693 token-exchange is cached per (jti, aud) — fan-out doesn't hammer Keycloak.",
    ]
    y = Inches(2.0)
    for p in points:
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.8), y + Inches(0.10),
                                    Inches(0.18), Inches(0.18))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = SAND
        bullet.line.fill.background()
        add_text(s, Inches(1.2), y, Inches(11.5), Inches(0.6), p,
                 size=16, color=WHITE)
        y += Inches(0.7)

    add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.5),
             "demo-a2a · architecture deck",
             size=12, color=SAND)


# ---------- main ------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_cover,
        slide_overview,
        slide_architecture,
        slide_keycloak,
        slide_hops,
        slide_flow_diagram,
        slide_orchestrator,
        slide_runtime,           # NEW: create_agent + SkillMiddleware
        slide_agent_internals,
        slide_token_exchange,
        slide_failure_paths,
        slide_add_agent,
        slide_techstack,
        slide_summary,
    ]
    for i, b in enumerate(builders, start=1):
        b(prs, i)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(out_dir, "demo-a2a-architecture.pptx")
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
